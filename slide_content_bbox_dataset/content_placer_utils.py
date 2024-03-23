import os
import random
import json
from pptx.util import Inches
import PIL.Image


"""
    load config : util 함수에 사용되는 configurations 
"""
with open("config.json", 'r') as config_file:
    config = json.load(config_file)
SLIDE_SIZE_PX = config['SLIDE_VARIANTS']['SLIDE_SIZE_PX']
SLIDE_BACKGROUND_IMG_DIR = config['SLIDE_VARIANTS']['SLIDE_BACKGROUND_IMG_DIR']
MAP_CONTENT_DIR = {
    "diagram": config['CONTENT_DIR']['diagram'],                   # list
    "natural_image": config['CONTENT_DIR']['natural_image'],       # list
    "text_paragraph": config['CONTENT_DIR']['text']['paragraph'],  # list
    "text_line": config['CONTENT_DIR']['text']['line'],            # list
    "text_short": config['CONTENT_DIR']['text']['short'],          # list
}
COCO_CATEGORIES = config['COCO_CATEGORIES']

"""
    Grid system 기반, Slide 내 Content 구성 및 배치 Variants 사전정의
    : 8x6 GRID 평면 내에서, 각각 [x0_gridpoint, y0_gridpoint, x1_gridpoint, y1_gridpoint]로 partition
"""
GRID_SIZE = (8, 6)
CONTENT_PLACEMENT_COMBINATION = [
    {
        "id": 1,
        "description": "1 main graphic, texts",
        "grid_placement_combination": [
            # horizontal middle
            {"main_graphic": [[1, 1, 5, 4]],
                "text_paragraph": [[5, 1, 8, 4]], "text_short": [[1, 4, 7, 6]]},
            {"main_graphic": [[2, 1, 6, 4]],  "text_short": [[1, 4, 7, 6]]},
            {"main_graphic": [[3, 1, 7, 4]],  "text_short": [[1, 4, 7, 6]]},
            # big
            {"main_graphic": [[1, 1, 6, 5]],  "text_short": [[1, 4, 7, 6]]},
            {"main_graphic": [[1, 1, 6, 6]]},
            # vertical middle
            {"main_graphic": [[1, 1, 4, 5]],
                "text_paragraph": [[4, 1, 7, 5]]},
            {"main_graphic": [[1, 1, 4, 5]],
                "text_paragraph": [[4, 1, 7, 3], [4, 3, 7, 5]]},
            {"main_graphic": [[1, 1, 4, 5]],
                "text_short": [[4, 1, 7, 2], [4, 2, 7, 3], [4, 3, 7, 4], [4, 4, 7, 5]]},
            {"main_graphic": [[4, 1, 7, 5]],
                "text_paragraph": [[1, 1, 4, 5]]},
            {"main_graphic": [[4, 1, 7, 5]],
                "text_paragraph": [[1, 1, 4, 3], [1, 3, 4, 5]]},
            {"main_graphic": [[4, 1, 7, 5]],
                "text_short": [[1, 1, 4, 2], [1, 2, 4, 3], [1, 3, 4, 4], [1, 4, 4, 5]]},
        ]
    },
    {
        "id": 2,
        "description": "2 main graphic, texts",
        "grid_placement_combination": [
            # 1 small 1 middle
            {"main_graphic": [[1, 1, 4, 3],  [4, 1, 7, 5]],
                "text_paragraph": [[1, 3, 4, 5]]},
            {"main_graphic": [[1, 1, 4, 3], [4, 1, 7, 5]],
                "text_short": [[1, 3, 4, 4], [1, 4, 4, 5]]},
            {"main_graphic": [[4, 1, 7, 3], [1, 1, 4, 5]],
                "text_paragraph": [[4, 3, 7, 5]]},
            {"main_graphic": [[4, 1, 7, 3], [1, 1, 4, 5]],
                "text_short": [[4, 3, 7, 4], [4, 4, 7, 5]]},
            # 2 middle
            {"main_graphic": [[1, 1, 4, 5], [4, 1, 7, 5]],
                "text_short": [[1, 5, 7, 6]]},
            # 2 small
            {"main_graphic": [[1, 1, 4, 3], [4, 1, 7, 3]],
                "text_paragraph": [[1, 3, 4, 5], [4, 3, 7, 5]]},
            {"main_graphic": [[1, 1, 4, 3], [4, 1, 7, 3]],
                "text_short": [[1, 3, 4, 4], [1, 4, 4, 5], [4, 3, 7, 4], [4, 4, 7, 5]]},
            {"main_graphic": [[1, 1, 4, 3], [4, 1, 7, 3]],
                "text_short": [[1, 3, 4, 4], [1, 4, 4, 5], [1, 5, 4, 6], [4, 3, 7, 4], [4, 4, 7, 5], [4, 5, 7, 6]]},
        ]
    },
    # {
    #     "id": 3,
    #     "description": "several sub graphics, texts",
    #     "grid_placement_combination": [
    #         # 3 small, free
    #         {"sub_graphic": [[1, 1, 3, 3], [5, 1, 7, 3], [3, 3, 6, 5]],
    #             "text_paragraph": [[3, 1, 5, 3], [1, 3, 3, 5], [6, 3, 7, 5]]},
    #         {"sub_graphic": [[1, 1, 3, 3], [5, 1, 7, 3], [3, 3, 6, 5]],
    #             "text_paragraph": [[3, 1, 5, 3], [1, 3, 3, 5], [6, 3, 7, 5]]},
    #         {"sub_graphic": [[1, 1, 3, 3], [5, 1, 7, 3], [3, 3, 6, 5]],
    #             "text_paragraph": [[3, 1, 5, 3], [1, 3, 3, 5], [6, 3, 7, 5]]}
    #         # 3 small, vertical

    #         # 2 middle, free

    #         # 1 middle, free

    #     ]
    # },
    # {
    #     "id": 4,
    #     "description": "no graphic, only texts",
    #     "grid_placement_combination": [
    #         # vertical

    #         # free

    #         # center

    #     ]
    # }
]
MAP_CONTENT_TYPE = {
    "main_graphic": ["diagram", "natural_image"],
    "sub_graphic": ["natural_image"],
    "text_paragraph": ["text_paragraph"],
    "text_line": ["text_line"],
    "text_short": ["text_short"],
}


"""
    Contents image path 모두 load 
"""


def find_imgs_in_deepest(root_dir):
    deepest_jpg_files = []
    for root, dirs, files in os.walk(root_dir):
        if not dirs:  # 가장 안쪽 폴더 순회
            for file in files:
                if file.endswith(('.png', '.jpg', '.jpeg')):
                    deepest_jpg_files.append(os.path.join(root, file))
    return deepest_jpg_files


CONTENT_JPGS_PATH = {
    "diagram": find_imgs_in_deepest(random.choice(MAP_CONTENT_DIR['diagram'])),
    "natural_image": find_imgs_in_deepest(random.choice(MAP_CONTENT_DIR['natural_image'])),
    "text_paragraph": find_imgs_in_deepest(random.choice(MAP_CONTENT_DIR['text_paragraph'])),
    "text_line": find_imgs_in_deepest(random.choice(MAP_CONTENT_DIR['text_line'])),
    "text_short": find_imgs_in_deepest(random.choice(MAP_CONTENT_DIR['text_short'])),
}

# CONTENT_JPGS_PATH에서 각 content의 이미지 갯수를 출력
for category, img_paths in CONTENT_JPGS_PATH.items():
    print(f"{category}: {len(img_paths)}")


"""
    util 함수 정의 
"""


def pixels_to_inches(pixels):
    # 단위 변환 : pixel -> inch
    return pixels / 96.0


def grid_to_inches(grid_x, grid_y):
    # 단위 변환 : grid -> inch
    grid_x_unit_px = SLIDE_SIZE_PX[0] / GRID_SIZE[0]
    grid_y_unit_px = SLIDE_SIZE_PX[1]/GRID_SIZE[1]

    grid_x_unit_in = pixels_to_inches(grid_x_unit_px)
    grid_y_unit_in = pixels_to_inches(grid_y_unit_px)

    grid_x_in = grid_x_unit_in * grid_x
    grid_y_in = grid_y_unit_in * grid_y

    return grid_x_in, grid_y_in


def make_slide(prs):
    # 슬라이드 크기/배경 지정하여 생성
    slide_width_px, slide_height_px = SLIDE_SIZE_PX[0], SLIDE_SIZE_PX[1]
    slide_width_in = pixels_to_inches(slide_width_px)
    slide_height_in = pixels_to_inches(slide_height_px)
    prs.slide_width = Inches(slide_width_in)
    prs.slide_height = Inches(slide_height_in)

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    bgimage_dir = SLIDE_BACKGROUND_IMG_DIR
    bgimage_files = [file for file in os.listdir(
        bgimage_dir) if file.endswith(('.png', '.jpg', '.jpeg'))]
    if bgimage_files:
        bgimage_file = random.choice(bgimage_files)
        bgimage_path = bgimage_dir + bgimage_file
        _ = slide.shapes.add_picture(
            bgimage_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    return slide


def _get_id_by_name(categories, key):
    for category in categories:
        if category['name'] == key:
            return category['id']
    return None


def _get_yolo_class_id(coco_categories, content_type_key):
    MAP_CATEGORIES = {
        "diagram": "diagram",
        "natural_image": "natural_image",
        "text_paragraph": "text",
        "text_line": "text",
        "text_short": "text",
    }
    class_id_yolo = _get_id_by_name(
        coco_categories, MAP_CATEGORIES[content_type_key])-1
    if class_id_yolo >= 0:
        return class_id_yolo
    return None


def select_layout_and_contents():
    # 가능한 레이아웃 중에서 선택 -> 레이아웃 각 구역에 들어갈 content 선택
    grid_layout_type = random.choice(CONTENT_PLACEMENT_COMBINATION)
    grid_layout = random.choice(grid_layout_type['grid_placement_combination'])
    print(grid_layout)
    content_type_list = list(grid_layout.keys())

    content_list = []
    for content_type in content_type_list:
        frame_grid_list = grid_layout[content_type]
        for frame_grid in frame_grid_list:
            # content 선택
            content_type_key = random.choice(MAP_CONTENT_TYPE[content_type])
            content_image_path = random.choice(
                CONTENT_JPGS_PATH[content_type_key])
            # frame partition 정보 저장
            class_id_yolo = _get_yolo_class_id(
                COCO_CATEGORIES, content_type_key)
            if content_image_path and (class_id_yolo != None):
                content_list.append(
                    {"class_id_yolo": class_id_yolo, "frame_grid": frame_grid, "content_image_path": content_image_path})
    # print(content_list)
    for item in content_list:
        print(item['content_image_path'])

    return content_list


def place_content_in_frame(slide, frame_grid, content_image_path):
    # 그리드 상의 박스 영역 frame_grid 안에서, 비율 유지하며 content_image_file를 fit하게 배치
    # + 그 과정에서 bbox 정보 자동 저장
    x0_grid, y0_grid, x1_grid, y1_grid = frame_grid
    frame_xcenter_in, frame_ycenter_in = grid_to_inches(
        (x0_grid+x1_grid)//2, (y0_grid+y1_grid)//2)
    frame_width_in, frame_height_in = grid_to_inches(
        x1_grid-x0_grid, y1_grid-y0_grid)
    frame_ratio = frame_width_in/frame_height_in
    with PIL.Image.open(content_image_path) as img:
        original_width, original_height = img.size
        image_ratio = original_width / original_height

    if image_ratio >= frame_ratio:
        # frame 가로 길이에 맞게, w/h 설정
        image_width_in = frame_width_in
        image_height_in = image_width_in*(original_height/original_width)
    else:
        # frame 세로 길이에 맞게, w/h 설정
        image_height_in = frame_height_in
        image_width_in = image_height_in*image_ratio
    # left/top 설정
    image_left_in = frame_xcenter_in-(image_width_in/2)
    image_top_in = frame_ycenter_in-(image_height_in/2)

    # slide에 content_image 배치
    _ = slide.shapes.add_picture(content_image_path,
                                 left=Inches(image_left_in),
                                 top=Inches(image_top_in),
                                 width=Inches(image_width_in),
                                 height=Inches(image_height_in))
    # print(image_height_in, Inches(image_height_in))

    # bbox 정보 반환
    slide_width_px, slide_height_px = SLIDE_SIZE_PX[0], SLIDE_SIZE_PX[1]
    slide_width_in = pixels_to_inches(slide_width_px)
    slide_height_in = pixels_to_inches(slide_height_px)
    x0, y0 = image_left_in/slide_width_in, image_top_in/slide_height_in,
    x1, y1 = (image_left_in+image_width_in)/slide_width_in, \
        (image_top_in+image_height_in)/slide_height_in
    bbox_yolo = [x0, y0, x1, y1]

    return slide, bbox_yolo
