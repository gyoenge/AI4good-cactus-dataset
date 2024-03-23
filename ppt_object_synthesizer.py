from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.util import Inches, Pt
import random
from pptx.shapes.autoshape import Shape
from pprint import pprint
import json


def k_modal_distribution(modes=[], weights=[]):
    value = sum([weight * random.gauss(mode) for (mode, weight) in zip(modes, weights)]) / sum(weights)
    return value if value > 0 else 0

def random_rgb_color():
    return RGBColor(random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))

def add_custom_manipulator(cls:Shape):
    def set(self, property, value):
        match property:
            case "line_type":
                self.line.dash_style = value
            case "left":
                self.left = value
            case "top":
                self.top = value
            case "width":
                self.width = value
            case "height":
                self.height = value
            case "rotation":
                self.rotation = value
            case "fill_fore_color":
                self.fill.solid()
                self.fill.fore_color.rgb = value
            case "line_color":
                self.line.color.rgb = value
            case "line_width":
                self.line.width = value

    def randomize(self):
        print(self.shape)
        random_properties = {
            # "shape_type": random.choices(
            #     [member.value for member in MSO_SHAPE.__members__], 
            #     [5 if member.name == "RECTANGLE" else 1 for member in MSO_SHAPE.__members__]
            # )[0], 
            "line_type": random.choice(
                [member.value for member in MSO_LINE.__members__ if member.name != "DASH_STYLE_MIXED"], 
            ),
            "left": Inches(random.uniform(0.5, 6.5)), 
            "top": Inches(random.uniform(0, 4)), 
            "width": Inches(random.uniform(1, 3)), 
            "height": Inches(random.uniform(1, 3)),  
            "rotation": random.choice([0, 45, 90, 135, 180, 225, 270, 315]),  # random.randint(0, 360),  # k_modal_distribution([0, 90, 180, 270], [3, 1, 3, 1]), 
            "fill_fore_color": random_rgb_color(), 
            "line_color": random_rgb_color(), 
            "line_width": Pt(random.randint(0, 5)),
        }
        for property, value in random_properties.items():
            self.set(property, value)
        setattr(cls, "properties", random_properties)
        return random_properties

    setattr(cls, "set", set)
    setattr(cls, "randomize", randomize)

    return cls

add_custom_manipulator(Shape)

class PPT:
    def __init__(self):
        self.prs = Presentation()

        self.prs.slide_width = Inches(8)
        self.prs.slide_height = Inches(6)

        self.prs.core_properties.author = "Cactus"
        self.prs.core_properties.category = "Image Dataset"
        self.prs.core_properties.comments = "This is a ppt file of randomly synthesized ppt shapes. It will be the dataset for training custom ppt object detection model. "
        self.prs.core_properties.title = "PPT Object Dataset"

        self.prs.core_properties.version = "0.0"

        self.shapes = dict((member.name, member.value) for member in MSO_SHAPE.__members__)
        shape_white_list = ["CLOUD", "CROSS", "LEFT_ARROW", "UP_ARROW", "RIGHT_ARROW", "DOWN_ARROW" "MATH_PLUS", "MATH_NOT_EQUAL", "RECTANGLE", "SIMLEY_FACE", "TRAPEZOID"]
        self.shapes = dict(shape for shape in self.shapes.items() if shape[0] in shape_white_list)


    def add_blank_slide(self):
        self.prs.slides.add_slide(self.prs.slide_layouts.get_by_name("Blank"))
        
    def generate_random_shape(self):
        shape_type = random.choice(list(self.shapes.items()))
        self.prs.slides[-1].shapes.add_shape(
            shape_type[1], 
            Inches(1), Inches(1), Inches(1),  Inches(1)
        )

        shape = self.prs.slides[-1].shapes[-1]
        setattr(shape, "shape", shape_type[0])

        property = shape.randomize()
        property["shape_type"] = shape_type[0]
        
        return property
    
    def _generate_random_rgb_color(self):
        return RGBColor(random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))

    def save(self):
        self.prs.save('ppt_object_dataset.pptx')

def synthesize_data(slides=10, shape_per_slide=4):
    ppt = PPT()
    answer = []
    for slide in range(slides):
        ppt.add_blank_slide()
        for shape in range(shape_per_slide):
            answer.append(ppt.generate_random_shape())
    ppt.save()
    with open('answer.json', 'w') as f:
        json.dump(answer, f, indent=4)

if __name__ == '__main__':
    synthesize_data()